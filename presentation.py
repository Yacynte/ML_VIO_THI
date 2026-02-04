from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Bullet layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = content_text_list[0] # First bullet
        
        for item in content_text_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "IMU-Based Pose Reconstruction\nusing Error-State EKF"
    slide.placeholders[1].text = "Analysis of KITTI Drive 2011_09_30_0020"

    # --- Slide 2: Objective ---
    add_slide("Objective & Data Source", [
        "Goal: Reconstruct 3D vehicle trajectory using only IMU data.",
        "Dataset: KITTI Odometry (2011_09_30_drive_0020_sync).",
        "Input Data:",
        "   - Accelerometer (ax, ay, az)",
        "   - Gyroscope (wx, wy, wz)",
        "   - Timestamps (converted to delta_t)"
    ])

    # --- Slide 3: Algorithm ---
    add_slide("Methodology: Error-State EKF", [
        "Architecture: Nominal State (Kinematics) + Error State (Noise).",
        "State Vector:",
        "   - Position (p), Velocity (v), Orientation (q)",
        "   - Biases (b_acc, b_gyro)",
        "Prediction Loop:",
        "   1. Integrate Angular Rate -> Update Quaternion",
        "   2. Rotate Accel to World Frame",
        "   3. Integrate Accel -> Velocity -> Position"
    ])

    # --- Slide 4: Physics ---
    add_slide("Key Challenge: Gravity Compensation", [
        "IMU measures Specific Force, not acceleration.",
        "We must subtract gravity in the World Frame:",
        "   a_world = (R * f_body) - G_world",
        "Crucial Dependency:",
        "   If Orientation (R) is wrong, Gravity is subtracted incorrectly.",
        "   This leaves 'leakage' acceleration on X/Y axes."
    ])

    # --- Slide 5: The Problem ---
    add_slide("The 'Straight Line' Drift Issue", [
        "Observation: Trajectory was a straight line, missing the 180-degree turn.",
        "Root Cause: Initialization Error.",
        "   - Initialized with q=[0,0,0,1] (Identity).",
        "   - Failed to account for initial vehicle Pitch/Roll.",
        "Result: A massive constant false acceleration masked the true turning motion."
    ])

    # --- Slide 6: The Solution ---
    add_slide("Solution: Static Initial Alignment", [
        "Method: Calculate q0 using first 50 static IMU frames.",
        "Roll (phi) = atan2(-ay, az)",
        "Pitch (theta) = atan2(ax, sqrt(ay^2 + az^2))",
        "Yaw (psi) = Obtained from OXTS heading.",
        "Outcome: Corrected the Gravity Vector alignment, revealing true motion."
    ])

    # --- Slide 7: Results ---
    add_slide("Results & Conclusion", [
        "[PLACEHOLDER: Insert your trajectory plot here]",
        "   - Blue line: EKF Prediction",
        "   - Orange line: Ground Truth (optional)",
        "Conclusion:",
        "   - Successfully implemented ES-EKF prediction step.",
        "   - Validated that correct Attitude Initialization is critical for INS accuracy."
    ])

    # Save
    prs.save('IMU_Pose_Reconstruction.pptx')
    print("Presentation saved as 'IMU_Pose_Reconstruction.pptx'")

if __name__ == "__main__":
    create_presentation()